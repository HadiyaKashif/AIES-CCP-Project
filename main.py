import os
from io import BytesIO
import io
import tempfile
import streamlit as st
import pytesseract
import requests
from bs4 import BeautifulSoup
from PIL import Image
from pptx import Presentation
from docx import Document
from dotenv import load_dotenv
from pinecone import Pinecone, ServerlessSpec
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_pinecone import PineconeVectorStore
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.runnables import RunnablePassthrough, RunnableLambda
from langchain_core.output_parsers import StrOutputParser
from langchain_google_genai import ChatGoogleGenerativeAI, GoogleGenerativeAIEmbeddings
from langchain_core.documents import Document as LangchainDocument
from langchain_google_community import GoogleSearchAPIWrapper
from langchain_core.tools import Tool
from collections import defaultdict
from typing import List
from fpdf import FPDF
import time
import json
import datetime
from streamlit_quill import st_quill
from io import BytesIO

st.set_page_config(page_title="Advanced Gemini RAG System", page_icon="🤖")

# Load .env secrets
load_dotenv()
PINECONE_API_KEY = os.getenv("PINECONE_API_KEY") or st.secrets["PINECONE_API_KEY"]
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY") or st.secrets["GEMINI_API_KEY"]
PINECONE_INDEX_NAME = "rag-advanced"
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY") or st.secrets["GOOGLE_API_KEY"]
GOOGLE_CSE_ID = os.getenv("GOOGLE_CSE_ID") or st.secrets["GOOGLE_CSE_ID"]

if GOOGLE_API_KEY and GOOGLE_CSE_ID:
    search = GoogleSearchAPIWrapper(
        google_api_key=GOOGLE_API_KEY,
        google_cse_id=GOOGLE_CSE_ID
    )
    web_search_tool = Tool(
        name="Google Search",
        description="Search Google for recent results",
        func=lambda query: search.run(query, num_results=3)  # Pass num_results here
    )
else:
    web_search_tool = None

def search_google(query, num_results=3):
    if not web_search_tool:
        return None, "Web search is not configured (missing API keys)"
    
    try:
        # Get raw JSON results
        raw_results = search.results(query, num_results)
        
        # Format the results
        formatted = []
        for i, result in enumerate(raw_results, 1):
            formatted.append(
                f"🔍 **Result {i}**\n"
                f"📌 {result.get('title', 'No title')}\n"
                f"🔗 {result.get('link', 'No link')}\n"
                f"📝 {result.get('snippet', 'No description')}\n"
            )
        return raw_results, "\n\n".join(formatted) if formatted else "No results found"
    except Exception as e:
        return None, f"Web search error: {str(e)}"

# Initialize Pinecone and Gemini
pc = Pinecone(api_key=PINECONE_API_KEY)
gemini_embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001", google_api_key=GEMINI_API_KEY)
gemini_llm = ChatGoogleGenerativeAI(model="gemini-2.0-flash", google_api_key=GEMINI_API_KEY, temperature=0.3)

if "messages" not in st.session_state:
    st.session_state.messages = []
if "processed" not in st.session_state:
    st.session_state.processed = False
if "vector_store" not in st.session_state:
    st.session_state.vector_store = None
if "rich_notes" not in st.session_state:
    st.session_state.rich_notes = ""
if "show_notes" not in st.session_state:
    st.session_state.show_notes = False 

# Chat Export Functions
def generate_chat_pdf(chat_history):
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    
    # Use standard PDF core font instead of Arial
    pdf.set_font("helvetica", size=12)
    
    for message in chat_history:
        role = message["role"].capitalize()
        content = message["content"]
        
        # Handle encoding and special characters
        try:
            content = content.encode('latin-1', 'replace').decode('latin-1')
        except:
            content = content.encode('utf-8', 'replace').decode('latin-1', 'replace')
        
        # Add message to PDF
        pdf.multi_cell(0, 10, f"{role}: {content}\n")
    
    # Create bytes buffer (updated method)
    pdf_buffer = BytesIO()
    pdf_output = pdf.output()  # Removed deprecated 'dest' parameter
    pdf_buffer.write(pdf_output)
    pdf_buffer.seek(0)
    return pdf_buffer

def export_chat_history(history, fmt):
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    base_name = f"chat_export_{timestamp}"
    
    if fmt == "json":
        return json.dumps(history, indent=2), "application/json", f"{base_name}.json"
    elif fmt == "txt":
        text = "\n".join([f"{msg['role'].capitalize()}: {msg['content']}" for msg in history])
        return text, "text/plain", f"{base_name}.txt"
    elif fmt == "pdf":
        pdf_bytes = generate_chat_pdf(history)
        return pdf_bytes, "application/pdf", f"{base_name}.pdf"

# Floating Notes Button
float_style = """
    <style>
        div[data-testid="stFloatNotes"] {
            position: fixed;
            top: 15px;
            right: 20px;
            z-index: 9999;
        }
        div[data-testid="stFloatNotes"] button {
            background-color: #f5f5f5;
            border-radius: 30px;
            padding: 0.5rem 0.8rem;
            border: 1px solid #ccc;
            box-shadow: 1px 1px 5px rgba(0,0,0,0.1);
        }
    </style>
""" 
st.markdown(float_style, unsafe_allow_html=True)
with st.container():
    st.markdown('<div data-testid="stFloatNotes">', unsafe_allow_html=True)
    if st.button("Notes + ", key="toggle_notes_btn"):
        st.session_state.show_notes = not st.session_state.get("show_notes", False)
    st.markdown('</div>', unsafe_allow_html=True)

def initialize_pinecone():
    try:
        existing_indexes = pc.list_indexes().names()
        
        if PINECONE_INDEX_NAME not in existing_indexes:
            st.warning("Creating new Pinecone index... (may take 1-2 minutes)")
            pc.create_index(
                name=PINECONE_INDEX_NAME,
                dimension=768,
                metric="cosine",
                spec=ServerlessSpec(cloud="aws", region="us-east-1")
            )
            
            # Wait for index to be ready
            with st.spinner("Waiting for index to be ready..."):
                while True:
                    try:
                        desc = pc.describe_index(PINECONE_INDEX_NAME)
                        if desc.status['ready']:
                            break
                        time.sleep(5)
                    except Exception as e:
                        time.sleep(5)
            
            st.success("Index created and ready to use!")
        else:
            # Check readiness for existing index
            with st.spinner("Checking index status..."):
                while True:
                    try:
                        desc = pc.describe_index(PINECONE_INDEX_NAME)
                        if desc.status['ready']:
                            break
                        time.sleep(5)
                    except Exception as e:
                        time.sleep(5)

            st.info(f"Using existing index: {PINECONE_INDEX_NAME}")

            if "cleared_on_startup" not in st.session_state:
                clear_index()
                st.session_state.cleared_on_startup = True
            
    except Exception as e:
        st.error(f"Pinecone initialization failed: {str(e)}")
        st.stop()

def clear_index():
    try:
        index = pc.Index(PINECONE_INDEX_NAME)
        try:
            index.delete(delete_all=True, namespace="")
        except Exception as ns_error:
            pass
        st.session_state.vector_store = None
        st.session_state.processed = False
        st.success("All documents cleared successfully!")
    except Exception as e:
        st.error(f"Error clearing index: {str(e)}")

# Website Scraper
def scrape_website(url):
    try:
        res = requests.get(url, headers={"User-Agent": "Mozilla/5.0"})
        soup = BeautifulSoup(res.text, "html.parser")
        return "\n".join([p.get_text() for p in soup.find_all("p")]) or "No relevant text found."
    except Exception as e:
        return f"Error scraping: {str(e)}"

# Text Chunking + Vector Store Push
def process_text_data(text_data, source):
    if text_data:
        doc = LangchainDocument(page_content=text_data, metadata={"source": source})
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        splits = text_splitter.split_documents([doc])
        try:
            if st.session_state.vector_store is None:
                st.session_state.vector_store = PineconeVectorStore.from_documents(
                    splits, gemini_embeddings, index_name=PINECONE_INDEX_NAME, namespace="default")
            else:
                st.session_state.vector_store.add_documents(splits)
            st.success(f"✅ Processed {len(splits)} chunks from {source}")
            st.session_state.processed = True
        except Exception as e:
            st.error(f"Vector store error: {str(e)}")

# File Processor
def process_documents(uploaded_files):
    for file in uploaded_files:
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file.name.split('.')[-1]}") as tmp:
            tmp.write(file.getvalue())
            tmp_path = tmp.name
        try:
            ext = file.name.split('.')[-1].lower()
            text = ""
            if ext == "pdf":
                loader = PyPDFLoader(tmp_path)
                pages = loader.load_and_split()
                text = "\n".join([p.page_content for p in pages])
            elif ext == "pptx":
                prs = Presentation(tmp_path)
                text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
            elif ext == "docx":
                doc = Document(tmp_path)
                text = "\n".join([p.text.strip() for p in doc.paragraphs if p.text.strip()])
            elif ext in ["png", "jpg", "jpeg"]:
                text = pytesseract.image_to_string(Image.open(tmp_path))
            else:
                st.error(f"Unsupported file type: {file.name}")
                continue
            process_text_data(text, file.name)
        except Exception as e:
            st.error(f"Processing error for {file.name}: {str(e)}")
        finally:
            os.unlink(tmp_path)

# Prompt Helpers
def generate_query_variations(q: str) -> List[str]:
    prompt = f"Generate 3 different rephrasings of this query:\nOriginal: {q}\n\nVariations:\n1. "
    return [q] + [v.strip() for v in gemini_llm.invoke(prompt).content.split("\n") if v.strip()][:3]

def generate_sub_questions(q: str) -> List[str]:
    prompt = f"Break this question into 2-3 standalone sub-questions:\nOriginal: {q}\n\nSub-questions:\n1. "
    return [v.strip() for v in gemini_llm.invoke(prompt).content.split("\n") if v.strip()][:3]

def generate_reasoning_steps(q: str) -> List[str]:
    prompt = f"To answer this, list steps to retrieve required info:\nQuestion: {q}\n\nSteps:\n1. "
    return [v.strip() for v in gemini_llm.invoke(prompt).content.split("\n") if v.strip()][:3]

def reciprocal_rank_fusion(results: List[List[LangchainDocument]], k=60) -> List[LangchainDocument]:
    scores = defaultdict(float)
    for docs in results:
        for rank, doc in enumerate(docs, 1):
            doc_id = doc.page_content
            scores[doc_id] += 1.0 / (rank + k)
    sorted_docs = sorted(scores.items(), key=lambda x: x[1], reverse=True)
    return [LangchainDocument(page_content=doc_id, metadata={}) for doc_id, _ in sorted_docs[:4]]

# UI
st.title("🤖 Advanced Gemini RAG System")
st.caption("Now with RAG-Fusion, Multi-Query, Decomposition + Notes")

initialize_pinecone()

with st.sidebar:
    st.header("📂 Data Management")
    uploaded_files = st.file_uploader("Upload files", type=["pdf", "pptx", "docx", "png", "jpg", "jpeg"], accept_multiple_files=True)
    url = st.text_input("Or enter website URL")
    col1, col2 = st.columns(2)
    if col1.button("Process Files", disabled=not uploaded_files):
        with st.spinner("Processing files..."):
            process_documents(uploaded_files)
    if col2.button("Scrape Site", disabled=not url):
        with st.spinner("Scraping site..."):
            process_text_data(scrape_website(url), url)
    if st.button("🧹 Clear All Data", type="secondary"):
        clear_index()
    
    # Added Chat Export UI
    st.divider()
    st.header("💾 Chat Export")
    export_format = st.radio("Format", options=["json", "txt", "pdf"])
    if st.button("Export Chat History"):
        if st.session_state.messages:
            data, mime, fname = export_chat_history(st.session_state.messages, export_format)
            st.download_button(
                label=f"Download {export_format.upper()}",
                data=data,
                file_name=fname,
                mime=mime
            )
        else:
            st.warning("No chat history to export!")

# Chat interface
st.divider()
st.header("🌐 Web Search Fallback")
use_web_search = st.checkbox("Search from web if no document matches", 
                           disabled=not (GOOGLE_API_KEY and GOOGLE_CSE_ID))

if use_web_search and not (GOOGLE_API_KEY and GOOGLE_CSE_ID):
    st.warning("Please set GOOGLE_API_KEY and GOOGLE_CSE_ID in .env to enable web search")

for msg in st.session_state.messages:
    with st.chat_message(msg["role"]):
        st.markdown(msg["content"])

# In your chat interface section, replace the current if/else logic with this:

if prompt := st.chat_input("Ask about your documents..."):
    st.session_state.messages.append({"role": "user", "content": prompt})
    with st.chat_message("user"):
        st.markdown(prompt)
    with st.chat_message("assistant"):
        if not st.session_state.processed and not use_web_search:
            st.warning("Please process or scrape some documents first.")
        else:
            with st.spinner("Thinking..."):
                try:
                    response = ""
                    document_answer_found = False
                    
                    # First try to get answer from documents if processed
                    if st.session_state.processed:
                        retriever = st.session_state.vector_store.as_retriever(search_kwargs={"k": 4})
                        queries = generate_query_variations(prompt) + generate_sub_questions(prompt) + generate_reasoning_steps(prompt)
                        all_docs = [retriever.invoke(q) for q in queries]
                        fused_docs = reciprocal_rank_fusion(all_docs)
                        
                        if fused_docs:
                            context = "\n\n".join([f"📄 Source {i+1}:\n{doc.page_content}" for i, doc in enumerate(fused_docs)])
                            template = """You are a helpful AI assistant. Answer the question based only on the context below.
- If you cannot answer from the context, say EXACTLY: "I couldn't find this information in the documents."
- Use bullet points for lists, steps, or comparisons.
- Do NOT include any headings or intro text.
- Stick strictly to the provided context.

Context:
{context}

Question: {question}

Answer:"""
                            prompt_template = ChatPromptTemplate.from_template(template)
                            rag_chain = (
                                {"context": RunnableLambda(lambda x: context), "question": RunnablePassthrough()}
                                | prompt_template
                                | gemini_llm
                                | StrOutputParser()
                            )
                            answer = rag_chain.invoke(prompt)
                            
                            # Check if the answer indicates missing information
                            if "couldn't find this information in the documents" not in answer.lower():
                                response = f"📄 Document Answer:\n\n{answer}"
                                document_answer_found = True
                    
                    # Handle cases where answer wasn't found
                    if not document_answer_found:
                        if use_web_search:
                            raw_web_results, formatted_web_results = search_google(prompt)
                            if raw_web_results and "error" not in str(formatted_web_results).lower():
                                # Process web results with Gemini
                                web_template = """Analyze these web search results and provide a comprehensive answer:
                                - Combine information from multiple sources if needed
                                - Always include source references like [1], [2] with corresponding links
                                - Keep the answer concise but informative
                                
                                Search Results:
                                {results}
                                
                                Question: {question}
                                
                                Answer with inline citations:"""
                                web_prompt = ChatPromptTemplate.from_template(web_template)
                                web_chain = (
                                    {"results": RunnablePassthrough(), "question": RunnablePassthrough()}
                                    | web_prompt
                                    | gemini_llm
                                    | StrOutputParser()
                                )
                                
                                processed_results = web_chain.invoke({
                                    "results": formatted_web_results, 
                                    "question": prompt
                                })
                                
                                # Add numbered source list
                                sources_section = "\n\n### References:\n" + "\n".join(
                                    f"{i+1}. [{res.get('title', 'Source')}]({res.get('link', '')})"
                                    for i, res in enumerate(raw_web_results)
                                )
                                
                                response = (
                                    f"🌐 Web Answer (since not found in documents):\n\n"
                                    f"{processed_results}\n"
                                    f"{sources_section}"
                                )
                            else:
                                response = "⚠️ Couldn't find relevant information in documents or through web search."
                        else:
                            response = (
                                "I couldn't find this information in your documents. "
                                "You can enable web search in the sidebar to search online for answers."
                            )
                    
                    st.markdown(response)
                    st.session_state.messages.append({"role": "assistant", "content": response})
                    
                except Exception as e:
                    st.error(f"Error generating answer: {str(e)}")

# Notes Section
if st.session_state.show_notes:
    with st.expander("📝 Notes Editor", expanded=True):
        temp_editor_content = st_quill(
            key="quill_editor",
            value=st.session_state.rich_notes,
            placeholder="Write your notes here...",
            html=True
        )

        if st.button("💾 Save Notes"):
            st.session_state.rich_notes = temp_editor_content
            st.success("Notes saved!")

    with st.expander("🔍 Preview Notes", expanded=False):
        st.markdown("### Rendered Notes:")
        st.markdown(st.session_state.rich_notes or "*No notes yet.*", unsafe_allow_html=True)

    # PDF Download Section
    if st.button("💾 Download Notes as PDF"):
        if st.session_state.rich_notes:
            soup = BeautifulSoup(st.session_state.rich_notes, "html.parser")
            clean_text = soup.get_text()

            pdf = FPDF()
            pdf.add_page()
            pdf.set_auto_page_break(auto=True, margin=15)
            pdf.set_font("Arial", size=12)

            for line in clean_text.split("\n"):
                pdf.multi_cell(0, 10, line)

            pdf_buffer = io.BytesIO()
            pdf_string = pdf.output(dest='S').encode('latin1')
            pdf_buffer = BytesIO()
            pdf_buffer.write(pdf_string)
            pdf_buffer.seek(0)

            st.download_button(
                label="📄 Confirm Download",
                data=pdf_buffer,
                file_name="my_notes.pdf",
                mime="application/pdf"
            )
        else:
            st.warning("No notes to download!, Save first")