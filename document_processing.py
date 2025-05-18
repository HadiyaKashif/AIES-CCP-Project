import os
import tempfile
import pytesseract
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document as LangchainDocument
from pptx import Presentation
from docx import Document
from PIL import Image
import requests
from bs4 import BeautifulSoup
from flask import flash, session
from vector_store import get_vector_store

def scrape_website(url):
    try:
        res = requests.get(url, headers={"User-Agent": "Mozilla/5.0"})
        soup = BeautifulSoup(res.text, "html.parser")
        return "\n".join([p.get_text() for p in soup.find_all("p")]) or "No relevant text found."
    except Exception as e:
        return f"Error scraping: {str(e)}"

def process_text_data(text_data, source):
    if text_data:
        doc = LangchainDocument(page_content=text_data, metadata={"source": source})
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        splits = text_splitter.split_documents([doc])
        try:
            vector_store = get_vector_store()
            vector_store.add_documents(splits)
            flash(f"Processed {len(splits)} chunks from {source}", "success")
            session['processed'] = True
            return True
        except Exception as e:
            flash(f"Vector store error: {str(e)}", "error")
            return False

def process_documents(uploaded_files):
    for file_path in uploaded_files:
        try:
            ext = os.path.splitext(file_path)[1].lower()[1:]  # Get extension without dot
            text = ""
            if ext == "pdf":
                try:
                    loader = PyPDFLoader(file_path)
                    pages = loader.load_and_split()
                    if not pages:
                        flash(f"Warning: No text content found in PDF file: {os.path.basename(file_path)}", "warning")
                        continue
                    text = "\n".join([p.page_content for p in pages])
                except ImportError as ie:
                    flash(f"PDF processing error: Missing required package. Please ensure 'pypdf' is installed. Error: {str(ie)}", "error")
                    continue
                except Exception as pe:
                    flash(f"PDF processing error for {os.path.basename(file_path)}: {str(pe)}", "error")
                    continue
            elif ext == "pptx":
                prs = Presentation(file_path)
                text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
            elif ext == "docx":
                doc = Document(file_path)
                text = "\n".join([p.text.strip() for p in doc.paragraphs if p.text.strip()])
            elif ext in ["png", "jpg", "jpeg"]:
                text = pytesseract.image_to_string(Image.open(file_path))
            else:
                flash(f"Unsupported file type: {os.path.basename(file_path)}", "error")
                continue
            
            if not text.strip():
                flash(f"Warning: No text content extracted from {os.path.basename(file_path)}", "warning")
                continue
                
            process_text_data(text, os.path.basename(file_path))
        except Exception as e:
            flash(f"Processing error for {os.path.basename(file_path)}: {str(e)}", "error")
            continue